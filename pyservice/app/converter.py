"""
PDF dönüştürme mantığı.

Bu modül OpenMF projesindeki converter.py temel alınarak uyarlanmıştır.
Strateji: önce LibreOffice denenir; PDF -> docx/pptx için LibreOffice
güvenilir olmadığından, uygun filtre yoksa PyMuPDF (fitz) + python-docx /
python-pptx ile saf-Python dönüşüme düşülür. Bu fallback her ortamda çalışır.
"""
import os
import subprocess
import tempfile
from pathlib import Path
from typing import Optional

SUPPORTED_CONVERSIONS = {
    'pdf': ['docx', 'pptx'],
    'docx': ['pdf'],
    'doc': ['pdf'],
    'pptx': ['pdf'],
    'ppt': ['pdf'],
    'xlsx': ['pdf'],
    'xls': ['pdf'],
    'html': ['pdf'],
    'htm': ['pdf'],
}

LIBREOFFICE_BIN = os.environ.get('LIBREOFFICE_BIN', 'soffice')


class ConversionError(Exception):
    pass


class UnsupportedConversionError(ConversionError):
    pass


class LibreOfficeNotFoundError(ConversionError):
    pass


class LibreOfficeNoFilterError(ConversionError):
    pass


def convert(input_path, target_format, output_dir=None, timeout=180, pre_ocr=False):
    input_path = Path(input_path)
    if not input_path.exists():
        raise FileNotFoundError(f"Girdi dosyası bulunamadı: {input_path}")

    input_ext = input_path.suffix.lower().lstrip('.')
    target_format = target_format.lower().lstrip('.')

    valid_targets = SUPPORTED_CONVERSIONS.get(input_ext)
    if valid_targets is None:
        raise UnsupportedConversionError(f"Desteklenmeyen girdi formatı: .{input_ext}")
    if target_format not in valid_targets:
        raise UnsupportedConversionError(
            f".{input_ext} -> .{target_format} dönüşümü desteklenmiyor."
        )

    if output_dir is None:
        output_dir = Path(tempfile.mkdtemp(prefix='pdf_tools_'))
    else:
        output_dir = Path(output_dir)
        output_dir.mkdir(parents=True, exist_ok=True)

    # Taranmış PDF'ler için OCR ön-işlemi (isteğe bağlı).
    if pre_ocr and input_ext == 'pdf':
        input_path = _ocr_pdf(input_path, output_dir)

    # PDF -> docx: önce pdf2docx (layout/tablo korumalı), başarısızsa PyMuPDF düz-metin.
    if input_ext == 'pdf' and target_format == 'docx':
        try:
            return _pdf2docx_convert(input_path, output_dir)
        except Exception as e:
            print(f"pdf2docx başarısız, düz-metin yöntemine düşülüyor: {e}")
            return _python_pdf_to_docx(input_path, output_dir)
    if input_ext == 'pdf' and target_format == 'pptx':
        return _python_pdf_to_pptx(input_path, output_dir)

    # Diğer tüm yönler (ofis -> pdf) LibreOffice ile güvenilir çalışıyor.
    return _libreoffice_convert(input_path, target_format, output_dir, timeout)


def _libreoffice_convert(input_path, target_format, output_dir, timeout=180):
    lo_filter = _get_lo_filter(input_path.suffix.lower().lstrip('.'), target_format)

    profile_dir = Path(output_dir) / 'lo_profile'
    profile_dir.mkdir(parents=True, exist_ok=True)

    cmd = [
        LIBREOFFICE_BIN, '--headless', '--norestore',
        f'-env:UserInstallation=file://{profile_dir}',
        '--convert-to', lo_filter,
        '--outdir', str(output_dir),
        str(input_path),
    ]

    try:
        result = subprocess.run(cmd, capture_output=True, text=True, timeout=timeout)
    except FileNotFoundError:
        raise LibreOfficeNotFoundError(f"LibreOffice bulunamadı: '{LIBREOFFICE_BIN}'")
    except subprocess.TimeoutExpired:
        raise ConversionError(f"Dönüşüm {timeout} saniyede tamamlanamadı.")

    if result.returncode != 0:
        err = (result.stderr or result.stdout or '').strip()
        raise ConversionError(f"LibreOffice dönüşümü başarısız: {err}")

    output_path = output_dir / (input_path.stem + '.' + target_format)
    if not output_path.exists():
        candidates = [p.name for p in output_dir.glob('*.*')]
        raise ConversionError(
            f"Beklenen çıktı '{output_path.name}' üretilmedi. Klasördeki dosyalar: {candidates}"
        )
    return str(output_path.resolve())


def merge_docx(input_paths, output_dir):
    """
    Birden fazla .docx dosyasını docxcompose ile sadık şekilde birleştirir:
    görseller, tablolar, stiller ve numaralandırma korunur.
    input_paths: sırayla birleştirilecek .docx yollarının listesi.
    """
    from docx import Document
    from docxcompose.composer import Composer

    if len(input_paths) < 2:
        raise ConversionError("Birleştirme için en az iki .docx dosyası gerekli.")

    output_dir = Path(output_dir)
    output_dir.mkdir(parents=True, exist_ok=True)

    master = Document(str(input_paths[0]))
    composer = Composer(master)

    for path in input_paths[1:]:
        doc = Document(str(path))
        composer.append(doc)

    output_path = output_dir / "birlestirilmis.docx"
    composer.save(str(output_path))

    if not output_path.exists() or output_path.stat().st_size == 0:
        raise ConversionError("Birleştirme çıktı üretemedi.")
    return str(output_path.resolve())


def _pdf2docx_convert(input_path, output_dir):
    """pdf2docx ile layout-korumalı dönüşüm: tablolar, sütunlar, fontlar korunur."""
    from pdf2docx import Converter

    output_path = Path(output_dir) / f"{input_path.stem}.docx"
    cv = Converter(str(input_path))
    try:
        # start/end verilmezse tüm sayfalar işlenir.
        cv.convert(str(output_path))
    finally:
        cv.close()

    if not output_path.exists() or output_path.stat().st_size == 0:
        raise ConversionError("pdf2docx çıktı üretemedi.")
    return str(output_path.resolve())


def _python_pdf_to_docx(input_path, output_dir):
    import fitz
    from docx import Document
    from docx.shared import Inches

    doc = Document()
    pdf = fitz.open(str(input_path))

    for page_num in range(pdf.page_count):
        page = pdf[page_num]
        text = page.get_text().strip()

        if text:
            doc.add_heading(f'Sayfa {page_num + 1}', level=2)
            for line in text.split('\n'):
                line = line.strip()
                if line:
                    doc.add_paragraph(line)
        else:
            doc.add_heading(f'Sayfa {page_num + 1} (görüntü sayfası)', level=2)

        for img_index, img in enumerate(page.get_images(full=True)):
            xref = img[0]
            base_image = pdf.extract_image(xref)
            image_path = output_dir / f"page{page_num + 1}_img{img_index}.{base_image['ext']}"
            with open(image_path, 'wb') as f:
                f.write(base_image['image'])
            try:
                doc.add_picture(str(image_path), width=Inches(5))
            except Exception:
                pass  # bozuk/desteklenmeyen gömülü görüntüyü atla

    pdf.close()
    output_path = output_dir / f"{input_path.stem}.docx"
    doc.save(str(output_path))
    return str(output_path.resolve())


def _python_pdf_to_pptx(input_path, output_dir):
    import fitz
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    pdf = fitz.open(str(input_path))

    for page_num in range(pdf.page_count):
        page = pdf[page_num]
        text = page.get_text().strip()

        slide = prs.slides.add_slide(prs.slide_layouts[6])
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
        title_box.text_frame.text = f'Sayfa {page_num + 1}'

        if text:
            content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(12), Inches(5))
            tf = content_box.text_frame
            tf.word_wrap = True
            for i, line in enumerate(text.split('\n')):
                line = line.strip()
                if not line:
                    continue
                if i == 0:
                    tf.text = line
                else:
                    tf.add_paragraph().text = line

        left = Inches(0.5)
        for img_index, img in enumerate(page.get_images(full=True)):
            xref = img[0]
            base_image = pdf.extract_image(xref)
            image_path = output_dir / f"page{page_num + 1}_img{img_index}.{base_image['ext']}"
            with open(image_path, 'wb') as f:
                f.write(base_image['image'])
            try:
                slide.shapes.add_picture(str(image_path), left, Inches(1.3), width=Inches(3))
                left += Inches(3.5)
            except Exception:
                pass

    pdf.close()
    output_path = output_dir / f"{input_path.stem}.pptx"
    prs.save(str(output_path))
    return str(output_path.resolve())


def _get_lo_filter(input_ext, target_format):
    mapping = {
        ('docx', 'pdf'): 'pdf:writer_pdf_Export',
        ('doc', 'pdf'): 'pdf:writer_pdf_Export',
        ('pptx', 'pdf'): 'pdf:impress_pdf_Export',
        ('ppt', 'pdf'): 'pdf:impress_pdf_Export',
        ('xlsx', 'pdf'): 'pdf:calc_pdf_Export',
        ('xls', 'pdf'): 'pdf:calc_pdf_Export',
        ('html', 'pdf'): 'pdf:writer_pdf_Export',
        ('htm', 'pdf'): 'pdf:writer_pdf_Export',
    }
    return mapping.get((input_ext, target_format), target_format)


def _ocr_pdf(input_path, output_dir):
    """ocrmypdf ile taranmış PDF'e metin katmanı ekler. Başarısızsa orijinali döndürür."""
    ocr_path = Path(output_dir) / f"{input_path.stem}_ocr.pdf"
    try:
        subprocess.run(
            ['ocrmypdf', '--skip-text', '--optimize', '0',
             '-l', os.environ.get('OCR_LANG', 'eng+tur'),
             str(input_path), str(ocr_path)],
            capture_output=True, text=True, timeout=300, check=True,
        )
        return ocr_path
    except Exception as e:
        print(f"OCR ön-işlemi başarısız (OCR'sız devam ediliyor): {e}")
        return input_path
